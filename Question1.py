import os
import csv
import mysql.connector
import pptx
import PyPDF2
from datetime import datetime
import re
from docx import Document
import fitz  # PyMuPDF

# ========================= DATABASE CONNECTION =========================
class MySQLStorage:
    def __init__(self, extractors, host="localhost", user="harshit", password="harshit", database="extracted_database5"):
        self.extractors = extractors if isinstance(extractors, list) else [extractors]
        self.host = host
        self.user = user
        self.password = password
        self.database = database
        self.conn = self._connect_db()
        self._initialize_db()

    def _connect_db(self):
        return mysql.connector.connect(
            host=self.host,
            user=self.user,
            password=self.password,
            database=self.database
        )

    def _initialize_db(self):
        cursor = self.conn.cursor()
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS ExtractedText (
                id INT AUTO_INCREMENT PRIMARY KEY,
                source VARCHAR(50),
                page_num INT,
                content LONGTEXT,
                font_name VARCHAR(100),
                font_size FLOAT,
                is_heading BOOLEAN,
                text_type VARCHAR(50),
                UNIQUE KEY (source, page_num, content(255))
            )
            """)
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS Tables (
                id INT AUTO_INCREMENT PRIMARY KEY,
                source VARCHAR(50),
                table_data LONGTEXT,
                UNIQUE KEY (source, table_data(255))
            )
        """)
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS Hyperlinks (
                id INT AUTO_INCREMENT PRIMARY KEY,
                source VARCHAR(50),
                page INT,
                text TEXT,
                url TEXT,
                UNIQUE KEY (source, page, url(255))
            )
        """)
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS Images (
                id INT AUTO_INCREMENT PRIMARY KEY,
                source VARCHAR(50),
                page INT,
                path TEXT,
                format VARCHAR(10),
                resolution VARCHAR(20),
                UNIQUE KEY (source, path(255))
            )
        """)
        self.conn.commit()
        cursor.close()

    def store_data(self):
        cursor = self.conn.cursor()
        for extractor in self.extractors:
            source_type = extractor.get_source_type()
            
            # Store text with metadata
            text_data = extractor.extract_text_with_metadata()
            for item in text_data:
                cursor.execute("""
                    INSERT IGNORE INTO ExtractedText 
                    (source, page_num, content, font_name, font_size, is_heading, text_type) 
                    VALUES (%s, %s, %s, %s, %s, %s, %s)
                """, (
                    source_type, 
                    item.get('page', 0), 
                    item.get('text', ''), 
                    item.get('font_name', 'Unknown'), 
                    item.get('font_size', 0.0), 
                    item.get('is_heading', False), 
                    item.get('text_type', 'Normal')
                ))
            
            # Store tables
            for table in extractor.extract_tables():
                table_string = "\n".join(["\t".join(row) for row in table["data"]])
                cursor.execute("INSERT IGNORE INTO Tables (source, table_data) VALUES (%s, %s)", 
                              (source_type, table_string))
            
            # Store hyperlinks
            for page, text, url in extractor.extract_hyperlinks():
                cursor.execute("INSERT IGNORE INTO Hyperlinks (source, page, text, url) VALUES (%s, %s, %s, %s)", 
                              (source_type, page, text, url))
            
            # Store image metadata
            for page, path, fmt, resolution in extractor.extract_images():
                cursor.execute("INSERT IGNORE INTO Images (source, page, path, format, resolution) VALUES (%s, %s, %s, %s, %s)", 
                              (source_type, page, path, fmt, resolution))
        
        self.conn.commit()
        cursor.close()
        self.conn.close()

# ========================= FILE LOADERS =========================
class PPTLoader:
    def __init__(self, file_path):
        self.file_path = file_path
        if not file_path.endswith('.pptx'):
            raise ValueError("Invalid PPT file")
    
    def load_file(self):
        return pptx.Presentation(self.file_path)

class PDFLoader:
    def __init__(self, file_path):
        self.file_path = file_path
        if not file_path.endswith('.pdf'):
            raise ValueError("Invalid PDF file")
    
    def load_file_pypdf(self):
        return PyPDF2.PdfReader(open(self.file_path, 'rb'))
    
    def load_file_pymupdf(self):
        return fitz.open(self.file_path)

import docx

class DOCXLoader:
    def __init__(self, file_path):
        self.file_path = file_path
        if not file_path.endswith('.docx'):
            raise ValueError("Invalid DOCX file")
    
    def load_docx(self):
        return docx.Document(self.file_path)


# ========================= DATA EXTRACTION =========================
class DataExtractor:
    def get_source_type(self):
        raise NotImplementedError("Subclasses must implement get_source_type method")
    
    def extract_text(self):
        raise NotImplementedError("Subclasses must implement extract_text method")
    
    def extract_text_with_metadata(self):
        raise NotImplementedError("Subclasses must implement extract_text_with_metadata method")
    
    def extract_hyperlinks(self):
        raise NotImplementedError("Subclasses must implement extract_hyperlinks method")
    
    def extract_images(self, output_dir):
        raise NotImplementedError("Subclasses must implement extract_images method")
    
    def extract_tables(self):
        raise NotImplementedError("Subclasses must implement extract_tables method")

class PPTDataExtractor(DataExtractor):
    def __init__(self, loader):
        self.loader = loader.load_file()
        self.file_path = loader.file_path

    def get_source_type(self):
        return "PPT"

    def extract_text(self):
        extracted_text = []
        for i, slide in enumerate(self.loader.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text.append(f"Page {i}: {shape.text}")
        return "\n".join(extracted_text)
    
    def extract_text_with_metadata(self):
        extracted_text = []
        for i, slide in enumerate(self.loader.slides, 1):
            # Get slide layout type
            layout_type = "Unknown"
            try:
                if hasattr(slide, "slide_layout") and hasattr(slide.slide_layout, "name"):
                    layout_type = slide.slide_layout.name
            except Exception:
                pass
            
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                
                # Check if shape has text
                if not shape.text_frame:
                    continue
                    
                # Safely determine if this is a placeholder and what type
                shape_type = "Normal"
                is_placeholder = False
                try:
                    # Check if shape has placeholder_format attribute and if it's valid
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        is_placeholder = True
                        if hasattr(shape, "placeholder_format") and hasattr(shape.placeholder_format, "type"):
                            ph_type = shape.placeholder_format.type
                            # 1 = title, 2 = body, 13 = text
                            if ph_type == 1:
                                shape_type = "Title"
                            elif ph_type == 2:
                                shape_type = "Body"
                            elif ph_type == 13:
                                shape_type = "Text"
                except Exception:
                    # If any error occurs, just use default value
                    pass
                
                # Process each paragraph in the shape
                for paragraph in shape.text_frame.paragraphs:
                    # Skip empty paragraphs
                    if not paragraph.text.strip():
                        continue
                        
                    # Get paragraph level safely
                    level = 0
                    try:
                        level = paragraph.level
                    except Exception:
                        pass
                    
                    # Determine if this is a heading based on context
                    is_heading = False
                    
                    # Method 1: Based on placeholder type
                    if shape_type == "Title":
                        is_heading = True
                    
                    # Method 2: Based on paragraph level and first paragraph in shape
                    elif level == 0 and paragraph == shape.text_frame.paragraphs[0]:
                        is_heading = True
                    
                    # Method 3: Based on font size comparison (if we can get font info)
                    try:
                        if paragraph.runs and len(paragraph.runs) > 0:
                            # Get average font size of the document
                            avg_size = self._get_average_font_size()
                            
                            # Get font size of this text
                            font = paragraph.runs[0].font
                            if hasattr(font, "size") and font.size:
                                font_size = font.size
                                if isinstance(font_size, (int, float)) and font_size > avg_size * 1.2:
                                    is_heading = True
                    except Exception:
                        pass
                    
                    # Get font info safely
                    font_name = "Default"
                    font_size = 0.0
                    if paragraph.runs:
                        try:
                            font = paragraph.runs[0].font
                            if hasattr(font, "name") and font.name:
                                font_name = font.name
                            
                            if hasattr(font, "size") and font.size:
                                # Handle different font size formats
                                if hasattr(font.size, "pt"):
                                    font_size = font.size.pt
                                else:
                                    try:
                                        # Try to convert to float
                                        font_size = float(font.size)
                                        # Convert from EMU to points if necessary
                                        if font_size > 1000:  # Likely in EMU
                                            font_size = font_size / 12700  # Convert EMU to points
                                    except (TypeError, ValueError):
                                        font_size = 0.0
                        except Exception:
                            pass
                    
                    # Add to extracted text items
                    extracted_text.append({
                        'page': i,
                        'text': paragraph.text,
                        'font_name': font_name,
                        'font_size': font_size,
                        'is_heading': is_heading,
                        'text_type': shape_type,
                        'layout': layout_type,
                        'level': level
                    })
        
        return extracted_text

    def _get_average_font_size(self):
        """Calculate the average font size in the presentation"""
        sizes = []
        for slide in self.loader.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                if hasattr(run.font, "size") and run.font.size:
                                    if hasattr(run.font.size, "pt"):
                                        sizes.append(run.font.size.pt)
                                    else:
                                        # Try to convert to float
                                        try:
                                            size = float(run.font.size)
                                            # Convert from EMU to points if necessary
                                            if size > 1000:  # Likely in EMU
                                                size = size / 12700  # Convert EMU to points
                                            sizes.append(size)
                                        except (TypeError, ValueError):
                                            pass
                            except Exception:
                                pass
        
        return sum(sizes) / len(sizes) if sizes else 12.0  # Default to 12pt if no valid sizes found

    def extract_hyperlinks(self):
        hyperlinks = []
        for i, slide in enumerate(self.loader.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                if hasattr(run, "hyperlink") and run.hyperlink and hasattr(run.hyperlink, "address") and run.hyperlink.address:
                                    hyperlinks.append((i, run.text, run.hyperlink.address))
                            except Exception:
                                pass
        return hyperlinks

    def extract_images(self, output_dir="output/ppt_images"):
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        images = []
        for i, slide in enumerate(self.loader.slides, 1):
            for j, shape in enumerate(slide.shapes):
                try:
                    if shape.shape_type == 13:  # Shape type 13 corresponds to images
                        image = shape.image
                        image_format = image.ext
                        width_px = int(shape.width / 9525)
                        height_px = int(shape.height / 9525)
                        resolution = f"{width_px}x{height_px}"
                        image_filename = f"{output_dir}/slide_{i}_img_{j+1}.{image_format}"
                        with open(image_filename, "wb") as f:
                            f.write(image.blob)
                        images.append((i, image_filename, image_format, resolution))
                except Exception:
                    pass
        return images

    def extract_tables(self):
        tables = []
        for i, slide in enumerate(self.loader.slides, 1):
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "table"):
                        table_data = []
                        row_count = len(shape.table.rows)
                        col_count = len(shape.table.columns)
                        for row in shape.table.rows:
                            table_data.append([cell.text for cell in row.cells])
                        tables.append({"page": i, "data": table_data, "size": f"{row_count}x{col_count}"})
                except Exception:
                    pass
        return tables

class PDFDataExtractor(DataExtractor):
    def __init__(self, loader):
        self.file_path = loader.file_path
        self.pdf_reader = loader.load_file_pypdf()
        self.pdf_doc = loader.load_file_pymupdf()

    def get_source_type(self):
        return "PDF"

    def extract_text(self):
        extracted_text = []
        for i, page in enumerate(self.pdf_reader.pages, 1):
            text = page.extract_text()
            if text:
                extracted_text.append(f"Page {i}: {text}")
        return "\n".join(extracted_text)
    
    def extract_text_with_metadata(self):
        extracted_text = []
        
        # Use PyMuPDF for detailed text extraction
        for i, page in enumerate(self.pdf_doc, 1):
            try:
                # Extract text blocks with style information
                blocks = page.get_text("dict")["blocks"]
                
                for block in blocks:
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line["spans"]:
                                text = span["text"].strip()
                                if not text:
                                    continue
                                    
                                font_name = span.get("font", "Unknown")
                                font_size = span.get("size", 0.0)
                                
                                # Determine if this is likely a heading based on font size
                                # Larger font typically indicates headings
                                avg_font_size = self._get_average_font_size()
                                is_heading = font_size > (avg_font_size * 1.2)  # 20% larger than average
                                
                                # Determine text type based on flags and style
                                text_type = "Normal"
                                flags = span.get("flags", 0)
                                
                                # Check font attributes based on flags
                                # Commonly: 1=bold, 2=italic, 4=underlined, etc.
                                if is_heading:
                                    text_type = "Heading"
                                elif flags & 1:  # Bold
                                    text_type = "Bold"
                                elif flags & 2:  # Italic
                                    text_type = "Italic"
                                elif flags & 4:  # Underlined
                                    text_type = "Underlined"
                                
                                # Create a metadata record
                                extracted_text.append({
                                    'page': i,
                                    'text': text,
                                    'font_name': font_name,
                                    'font_size': font_size,
                                    'is_heading': is_heading,
                                    'text_type': text_type,
                                    'flags': flags,
                                    'color': span.get("color", 0)
                                })
            except Exception as e:
                print(f"Warning: Error extracting text with metadata from page {i}: {e}")
                # Add a basic entry to avoid missing pages
                extracted_text.append({
                    'page': i,
                    'text': f"Error extracting text from page {i}",
                    'font_name': "Unknown",
                    'font_size': 0.0,
                    'is_heading': False,
                    'text_type': "Error",
                })
        
        return extracted_text
    
    def _get_average_font_size(self):
        # Calculate average font size across the document
        sizes = []
        for page in self.pdf_doc:
            try:
                blocks = page.get_text("dict")["blocks"]
                for block in blocks:
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line["spans"]:
                                if "size" in span and span["size"] > 0:
                                    sizes.append(span["size"])
            except Exception:
                pass
        
        return sum(sizes) / len(sizes) if sizes else 12.0  # Default to 12pt if no valid sizes found

    def extract_hyperlinks(self):
        hyperlinks = []
        for i, page in enumerate(self.pdf_doc, 1):
            try:
                links = page.get_links()
                for link in links:
                    if 'uri' in link:
                        # Get text near the link
                        rect = fitz.Rect(link['from'])
                        # Extend the rect by a small amount to capture text
                        rect.x0 -= 5
                        rect.y0 -= 5
                        rect.x1 += 5
                        rect.y1 += 5
                        
                        # Get words within the extended rectangle
                        words = page.get_text("words", clip=rect)
                        link_text = " ".join([word[4] for word in words]) if words else "Unknown link text"
                        hyperlinks.append((i, link_text, link['uri']))
            except Exception as e:
                print(f"Warning: Error extracting hyperlinks from page {i}: {e}")
        return hyperlinks

    def extract_images(self, output_dir="output/pdf_images"):
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        images = []
        for i, page in enumerate(self.pdf_doc, 1):
            try:
                image_list = page.get_images(full=True)
                for j, img in enumerate(image_list, 1):
                    try:
                        xref = img[0]
                        base_image = self.pdf_doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        image_filename = f"{output_dir}/page_{i}_img_{j}.{image_ext}"
                        with open(image_filename, "wb") as f:
                            f.write(image_bytes)
                        
                        # Get dimensions from the image if available
                        width = base_image.get("width", 0)
                        height = base_image.get("height", 0)
                        resolution = f"{width}x{height}" if width and height else "unknown"
                        images.append((i, image_filename, image_ext, resolution))
                    except Exception as e:
                        print(f"Warning: Error extracting image {j} from page {i}: {e}")
            except Exception as e:
                print(f"Warning: Error getting images from page {i}: {e}")
        return images

    def extract_tables(self):
        tables = []
        try:
            # Try to import pandas for table handling
            import pandas as pd
            
            # Simple table detection using PyMuPDF
            for i, page in enumerate(self.pdf_doc, 1):
                try:
                    # Look for potential tables by finding rectangles that might represent cells
                    rect_areas = []
                    for rect in page.search_for("   "):  # Search for multiple spaces as table indicator
                        rect_areas.append(rect)
                    
                    # Very simple table detection heuristic
                    # Get all text blocks
                    blocks = page.get_text("blocks")
                    
                    # Filter blocks that might be table cells (simple heuristic)
                    potential_table_cells = []
                    for block in blocks:
                        rect = fitz.Rect(block[:4])
                        text = block[4].strip()
                        if text and len(text) < 100:  # Likely a cell if short text
                            potential_table_cells.append((rect, text))
                    
                    # Skip if not enough potential cells
                    if len(potential_table_cells) < 4:
                        continue
                    
                    # Simple row detection (group by similar y-coordinates)
                    y_coords = {}
                    for rect, text in potential_table_cells:
                        y_key = round(rect.y0 / 10) * 10  # Group by 10-pixel increments
                        if y_key not in y_coords:
                            y_coords[y_key] = []
                        y_coords[y_key].append((rect, text))
                    
                    # If we have at least 2 rows, consider it a table
                    if len(y_coords) >= 2:
                        # Convert to table data format
                        table_data = []
                        for y_key in sorted(y_coords.keys()):
                            row = []
                            # Sort cells by x-coordinate
                            sorted_cells = sorted(y_coords[y_key], key=lambda c: c[0].x0)
                            for _, text in sorted_cells:
                                row.append(text)
                            if row:
                                table_data.append(row)
                        
                        # Add to tables list if we have data
                        if table_data and len(table_data) > 1:  # At least 2 rows
                            row_count = len(table_data)
                            col_count = max(len(row) for row in table_data)
                            tables.append({"page": i, "data": table_data, "size": f"{row_count}x{col_count}"})
                except Exception as e:
                    print(f"Warning: Error extracting tables from page {i}: {e}")
        
        except Exception as e:
            print(f"Warning: Table extraction failed: {e}")
            # Minimal fallback for table detection
            for i, page in enumerate(self.pdf_doc, 1):
                try:
                    text = page.get_text()
                    # Check for table-like patterns
                    if "\t" in text or "  " in text:
                        rows = text.split("\n")
                        table_data = [row.split("\t" if "\t" in row else "  ") for row in rows if row.strip()]
                        if table_data and sum(len(row) > 1 for row in table_data) > 3:  # At least 3 rows with multiple columns
                            row_count = len(table_data)
                            col_count = max(len(row) for row in table_data)
                            tables.append({"page": i, "data": table_data, "size": f"{row_count}x{col_count}"})
                except Exception as e:
                    print(f"Warning: Error using fallback table extraction for page {i}: {e}")
        
        return tables

# ========================= FILE EXPORT =========================
def save_text_and_metadata(extractor, output_dir="output", prefix=""):
    os.makedirs(output_dir, exist_ok=True)
    
    # Save simple text (for backward compatibility)
    with open(f"{output_dir}/{prefix}_text.txt", "w", encoding="utf-8") as text_file:
        text_file.write(extractor.extract_text())
    
    # Save detailed text with metadata
    with open(f"{output_dir}/{prefix}_text_metadata.csv", "w", newline="", encoding="utf-8") as csv_file:
        writer = csv.writer(csv_file)
        writer.writerow(["Page", "Text", "Font Name", "Font Size", "Is Heading", "Text Type"])
        
        for item in extractor.extract_text_with_metadata():
            writer.writerow([
                item.get('page', 0),
                item.get('text', ''),
                item.get('font_name', 'Unknown'),
                item.get('font_size', 0.0),
                "Yes" if item.get('is_heading', False) else "No",
                item.get('text_type', 'Normal')
            ])

def save_tables_and_metadata(extractor, output_dir="output", prefix=""):
    os.makedirs(output_dir, exist_ok=True)
    with open(f"{output_dir}/{prefix}_tables.csv", "w", newline="", encoding="utf-8") as csv_file:
        writer = csv.writer(csv_file)
        for table in extractor.extract_tables():
            writer.writerow([f"Page {table['page']} - Size: {table['size']}"])
            writer.writerows(table["data"])
            writer.writerow([])

def save_links_and_metadata(extractor, output_dir="output", prefix=""):
    os.makedirs(output_dir, exist_ok=True)
    with open(f"{output_dir}/{prefix}_hyperlinks.txt", "w", encoding="utf-8") as link_file:
        for page, text, url in extractor.extract_hyperlinks():
            link_file.write(f"Page {page}: {text} -> {url}\n")

def save_images_metadata(extractor, output_dir="output", prefix=""):
    img_output_dir = f"{output_dir}/{prefix}_images"
    os.makedirs(output_dir, exist_ok=True)
    with open(f"{output_dir}/{prefix}_images_metadata.txt", "w", encoding="utf-8") as image_file:
        for page, path, fmt, resolution in extractor.extract_images(img_output_dir):
            image_file.write(f"Page {page}: {path} (Format: {fmt}, Resolution: {resolution})\n")

import os
import csv
from PIL import Image
from io import BytesIO

class DOCXLoader:
    def __init__(self, file_path):
        self.file_path = file_path
        if not file_path.endswith('.docx'):
            raise ValueError("Invalid DOCX file")
    
    def load_docx(self):
        try:
            from docx import Document
            return Document(self.file_path)
        except Exception as e:
            raise Exception(f"Error loading DOCX file: {e}")

class DOCXDataExtractor(DataExtractor):
    def __init__(self, loader):
        self.file_path = loader.file_path
        self.doc = loader.load_docx()

    def get_source_type(self):
        return "DOCX"

    def extract_text(self):
        extracted_text = []
        # Extract text from paragraphs
        for para in self.doc.paragraphs:
            if para.text.strip():
                extracted_text.append(para.text.strip())
        
        # Extract text from tables
        for table in self.doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        if para.text.strip():
                            extracted_text.append(para.text.strip())
        
        return "\n".join(extracted_text)

    def extract_text_with_metadata(self):
        extracted_text = []
        page_num = 1  # DOCX doesn't have built-in page numbers, so we estimate
        
        # Calculate average font size for reference
        avg_font_size = self._get_average_font_size()
        
        # Process paragraphs
        for para_index, para in enumerate(self.doc.paragraphs):
            if not para.text.strip():
                continue
            
            # Check if it's a heading
            is_heading = para.style.name.startswith("Heading")
            
            # Get font information from the first run if available
            font_name = "Default"
            font_size = 0.0
            text_type = "Normal"
            
            if para.runs:
                try:
                    run = para.runs[0]
                    if hasattr(run, "font") and run.font:
                        if hasattr(run.font, "name") and run.font.name:
                            font_name = run.font.name
                        
                        # Properly extract font size
                        if hasattr(run.font, "size") and run.font.size:
                            # Convert Pt object to float if needed
                            if hasattr(run.font.size, "pt"):
                                font_size = run.font.size.pt
                            else:
                                try:
                                    # Try direct conversion
                                    font_size = float(run.font.size)
                                    # DOCX sometimes stores size in half-points
                                    if font_size > 100:  # Likely in twips (1/20 of a point)
                                        font_size = font_size / 20
                                except (TypeError, ValueError):
                                    # Fall back to style-based size estimation
                                    style_font_size = self._get_font_size_from_style(para.style)
                                    if style_font_size > 0:
                                        font_size = style_font_size
                except Exception as e:
                    print(f"Warning: Error getting font info: {e}")
            
            # If we still don't have font size, estimate from style
            if font_size == 0.0:
                font_size = self._get_font_size_from_style(para.style)
                
            # Determine text type based on style
            if is_heading:
                text_type = para.style.name  # e.g., "Heading 1"
            elif para.style.name != "Normal":
                text_type = para.style.name
            
            # Estimate page number (rough approximation)
            if para_index > 0 and para_index % 15 == 0:  # Assume ~15 paragraphs per page
                page_num += 1
            
            # Determine if this is a heading based on context if not already set
            if not is_heading:
                # Check if font size is significantly larger than average
                if font_size > (avg_font_size * 1.2):
                    is_heading = True
                    text_type = "Possible Heading"
            
            extracted_text.append({
                'page': page_num,
                'text': para.text.strip(),
                'font_name': font_name,
                'font_size': font_size,
                'is_heading': is_heading,
                'text_type': text_type,
                'style': para.style.name
            })
        
        # Process text in tables
        for table_index, table in enumerate(self.doc.tables):
            for row_index, row in enumerate(table.rows):
                for cell_index, cell in enumerate(row.cells):
                    for para in cell.paragraphs:
                        if not para.text.strip():
                            continue
                        
                        # Similar font extraction for table cells
                        font_name = "Default"
                        font_size = 0.0
                        text_type = "Table Cell"
                        
                        # Try to get cell style first
                        cell_style = None
                        if hasattr(cell, "paragraphs") and cell.paragraphs:
                            if hasattr(cell.paragraphs[0], "style"):
                                cell_style = cell.paragraphs[0].style
                        
                        # Extract font info from all runs in the paragraph, not just the first
                        if para.runs:
                            for run in para.runs:
                                try:
                                    if hasattr(run, "font") and run.font:
                                        # Get font name if we don't have one yet
                                        if font_name == "Default" and hasattr(run.font, "name") and run.font.name:
                                            font_name = run.font.name
                                        
                                        # Try to get font size if we don't have a valid one yet
                                        if font_size == 0.0 and hasattr(run.font, "size") and run.font.size:
                                            if hasattr(run.font.size, "pt"):
                                                font_size = run.font.size.pt
                                                break  # Once we have a valid size, we can stop
                                            else:
                                                try:
                                                    size = float(run.font.size)
                                                    # Convert from twips if necessary
                                                    if size > 100:  # Likely in twips
                                                        size = size / 20
                                                    font_size = size
                                                    break  # Once we have a valid size, we can stop
                                                except (TypeError, ValueError):
                                                    continue
                                except Exception:
                                    continue
                        
                        # If we still don't have font size, try to get it from paragraph style
                        if font_size == 0.0 and hasattr(para, "style") and para.style:
                            font_size = self._get_font_size_from_style(para.style)
                        
                        # If still no font size, try cell style
                        if font_size == 0.0 and cell_style:
                            font_size = self._get_font_size_from_style(cell_style)
                        
                        # Last resort: use average document font size or default
                        if font_size == 0.0:
                            font_size = avg_font_size
                        
                        # Get cell position info for metadata
                        position_info = f"Table {table_index+1}, Row {row_index+1}, Cell {cell_index+1}"
                        
                        # First row might be headers
                        is_heading = row_index == 0
                        if is_heading:
                            text_type = "Table Header"
                        
                        # Store the original paragraph style name
                        para_style_name = para.style.name if hasattr(para, "style") and para.style and hasattr(para.style, "name") else "Default"
                        
                        extracted_text.append({
                            'page': page_num,  # Use the current page estimate
                            'text': para.text.strip(),
                            'font_name': font_name,
                            'font_size': font_size,
                            'is_heading': is_heading,
                            'text_type': text_type,
                            'style': para_style_name  # Use the actual paragraph style name instead of position_info
                        })
        
        return extracted_text

    def _get_average_font_size(self):
        """Calculate the average font size in the document"""
        sizes = []
        for para in self.doc.paragraphs:
            for run in para.runs:
                try:
                    if hasattr(run.font, "size") and run.font.size:
                        if hasattr(run.font.size, "pt"):
                            sizes.append(run.font.size.pt)
                        else:
                            try:
                                size = float(run.font.size)
                                # Convert from twips if necessary
                                if size > 100:  # Likely in twips
                                    size = size / 20
                                sizes.append(size)
                            except (TypeError, ValueError):
                                pass
                except Exception:
                    pass
        
        # Also check tables
        for table in self.doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            try:
                                if hasattr(run.font, "size") and run.font.size:
                                    if hasattr(run.font.size, "pt"):
                                        sizes.append(run.font.size.pt)
                                    else:
                                        try:
                                            size = float(run.font.size)
                                            if size > 100:  # Likely in twips
                                                size = size / 20
                                            sizes.append(size)
                                        except (TypeError, ValueError):
                                            pass
                            except Exception:
                                pass
        
        return sum(sizes) / len(sizes) if sizes else 12.0  # Default to 12pt if no valid sizes found

    def _get_font_size_from_style(self, style):
        """Estimate font size based on style name"""
        if not style or not style.name:
            return 12.0  # Default size
        
        # Common heading sizes
        if style.name == "Heading 1":
            return 16.0
        elif style.name == "Heading 2":
            return 14.0
        elif style.name == "Heading 3":
            return 13.0
        elif style.name == "Heading 4":
            return 12.0
        elif style.name == "Title":
            return 18.0
        elif style.name == "Subtitle":
            return 16.0
        elif "Caption" in style.name:
            return 10.0
        elif "Footer" in style.name or "Header" in style.name:
            return 10.0
        else:
            return 12.0  # Default body text size
    def extract_hyperlinks(self):
        hyperlinks = []
        for para_index, para in enumerate(self.doc.paragraphs):
            # Extracting hyperlinks from paragraphs
            rels = para._element.xpath('.//w:hyperlink')
            for rel in rels:
                rel_id = rel.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rel_id and rel_id in para.part.rels:
                    url = para.part.rels[rel_id].target_ref
                    text_elements = rel.xpath('.//w:t')
                    text = ''.join([t.text for t in text_elements if t.text])
                    hyperlinks.append((para_index + 1, text or "Unnamed link", url))
        
        # Check tables as well
        for table_index, table in enumerate(self.doc.tables):
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        rels = para._element.xpath('.//w:hyperlink')
                        for rel in rels:
                            rel_id = rel.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                            if rel_id and rel_id in para.part.rels:
                                url = para.part.rels[rel_id].target_ref
                                text_elements = rel.xpath('.//w:t')
                                text = ''.join([t.text for t in text_elements if t.text])
                                hyperlinks.append((table_index + 1, f"Table {table_index + 1}: {text or 'Unnamed link'}", url))
        
        return hyperlinks

    def extract_images(self, output_dir="output/docx_images"):
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        images = []
        page_num = 1  # Approximate page number
        
        # Extract images using document relationships
        for i, rel in enumerate(self.doc.part.rels.values()):
            if "image" in rel.reltype:
                try:
                    # Get the image data
                    image_data = rel.target_part.blob
                    
                    # Determine image format from relationship target
                    image_ext = rel.target_ref.split(".")[-1].lower()
                    if image_ext not in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff']:
                        image_ext = 'png'  # Default to PNG if unknown extension
                    
                    # Create a unique filename
                    image_filename = f"{output_dir}/image_{i+1}.{image_ext}"
                    
                    # Save the image file
                    with open(image_filename, "wb") as f:
                        f.write(image_data)
                    
                    # Get image dimensions
                    try:
                        image = Image.open(BytesIO(image_data))
                        resolution = f"{image.width}x{image.height}"
                    except:
                        resolution = "unknown"
                    
                    # Approximate page number (increments every 3 images)
                    if i > 0 and i % 3 == 0:
                        page_num += 1
                    
                    images.append((page_num, image_filename, image_ext, resolution))
                except Exception as e:
                    print(f"Warning: Failed to extract image {i+1}: {e}")
        
        return images
    
    def extract_tables(self):
        tables = []
        
        # Process each table in the document
        for table_index, table in enumerate(self.doc.tables):
            try:
                table_data = []
                
                # Process each row
                for row in table.rows:
                    row_data = []
                    
                    # Process each cell in the row
                    for cell in row.cells:
                        # Combine all paragraphs in the cell
                        cell_text = "\n".join([p.text.strip() for p in cell.paragraphs if p.text.strip()])
                        row_data.append(cell_text)
                    
                    # Add the row data to the table
                    if any(row_data):  # Only add non-empty rows
                        table_data.append(row_data)
                
                # Calculate table dimensions
                row_count = len(table_data)
                col_count = max([len(row) for row in table_data]) if table_data else 0
                
                # Add to tables list with page approximation
                tables.append({
                    "page": 1 + (table_index // 2),  # Rough approximation: 2 tables per page
                    "data": table_data,
                    "size": f"{row_count}x{col_count}"
                })
            except Exception as e:
                print(f"Warning: Error extracting table {table_index+1}: {e}")
        
        return tables

def save_docx_text_and_metadata(extractor, output_dir="output", prefix="docx"):
    """Special function for saving DOCX extracted data with appropriate columns"""
    os.makedirs(output_dir, exist_ok=True)
    
    # Save plain text
    with open(f"{output_dir}/{prefix}_text.txt", "w", encoding="utf-8") as text_file:
        text_file.write(extractor.extract_text())
    
    # Save detailed metadata
    with open(f"{output_dir}/{prefix}_text_metadata.csv", "w", newline="", encoding="utf-8") as csv_file:
        writer = csv.writer(csv_file)
        writer.writerow(["Page", "Text", "Font Name", "Font Size", "Is Heading", "Text Type", "Style"])
        
        for item in extractor.extract_text_with_metadata():
            writer.writerow([
                item.get('page', 0),
                item.get('text', ''),
                item.get('font_name', 'Unknown'),
                item.get('font_size', 0.0),
                "Yes" if item.get('is_heading', False) else "No",
                item.get('text_type', 'Normal'),
                item.get('style', 'Default')
            ])
# ========================= MAIN =========================
if __name__ == "__main__":
    try:
        extractors = []
        
        # PPT Processing
        try:
            ppt_path = "project.pptx"
            ppt_loader = PPTLoader(ppt_path)
            ppt_extractor = PPTDataExtractor(ppt_loader)
            extractors.append(ppt_extractor)
            
            print("\n🔹 Extracting PPT Content...")
            save_text_and_metadata(ppt_extractor, prefix="ppt")
            save_tables_and_metadata(ppt_extractor, prefix="ppt")
            save_links_and_metadata(ppt_extractor, prefix="ppt")
            save_images_metadata(ppt_extractor, prefix="ppt")
            print("✅ PPT data successfully extracted!")
        except Exception as e:
            print(f"❌ Error processing PPT: {e}")
        
        # PDF Processing
        try:
            pdf_path = "report.pdf"
            pdf_loader = PDFLoader(pdf_path)
            pdf_extractor = PDFDataExtractor(pdf_loader)
            extractors.append(pdf_extractor)
            
            print("\n🔹 Extracting PDF Content...")
            save_text_and_metadata(pdf_extractor, prefix="pdf")
            save_tables_and_metadata(pdf_extractor, prefix="pdf")
            save_links_and_metadata(pdf_extractor, prefix="pdf")
            save_images_metadata(pdf_extractor, prefix="pdf")
            print("✅ PDF data successfully extracted!")
        except Exception as e:
            print(f"❌ Error processing PDF: {e}")
        
        # DOCX Processing
        try:
            docx_path = "Sample6.docx"
            docx_loader = DOCXLoader(docx_path)
            docx_extractor = DOCXDataExtractor(docx_loader)
            extractors.append(docx_extractor)
            
            print("\n🔹 Extracting DOCX Content...")
            save_text_and_metadata(docx_extractor, prefix="docx")
            save_tables_and_metadata(docx_extractor, prefix="docx")
            save_links_and_metadata(docx_extractor, prefix="docx")
            save_images_metadata(docx_extractor, prefix="docx")
            print("✅ DOCX data successfully extracted!")
        except Exception as e:
            print(f"❌ Error processing DOCX: {e}")
        
        # Store in MySQL (if we have any extractors)
        if extractors:
            print("\n🔹 Storing data in MySQL...")
            storage = MySQLStorage(extractors)
            storage.store_data()
            print("✅ Data successfully stored in MySQL!")
        
        print("\n✅ Processing completed!")
    
    except Exception as e:
        print(f"❌ Error in main process: {e}")
